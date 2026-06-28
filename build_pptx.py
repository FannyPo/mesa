from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK  = RGBColor(0x1a, 0x3a, 0x2a)
MID   = RGBColor(0x2d, 0x6a, 0x4f)
LIGHT = RGBColor(0x52, 0xb7, 0x88)
CREAM = RGBColor(0xfa, 0xf7, 0xf2)
CHAR  = RGBColor(0x1c, 0x1c, 0x1e)
WHITE = RGBColor(0xff, 0xff, 0xff)
SAGE_BG = RGBColor(0xe8, 0xf5, 0xee)
WARM_BG = RGBColor(0xec, 0xe8, 0xe2)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(sl, c):
    f = sl.background.fill; f.solid(); f.fore_color.rgb = c

def rect(sl, l, t, w, h, c):
    shp = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = c
    shp.line.color.rgb = c
    return shp

def tb(sl, text, l, t, w, h, size=17, bold=False, color=CHAR,
       align=PP_ALIGN.LEFT, font='Lato', italic=False):
    box = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return box, tf

def ap(tf, text, size=15, bold=False, color=CHAR, font='Lato',
       align=PP_ALIGN.LEFT, before=7, italic=False):
    p = tf.add_paragraph(); p.alignment = align; p.space_before = Pt(before)
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic; r.font.color.rgb = color

def hdr(sl, title):
    bg(sl, CREAM)
    rect(sl, 0, 0, 13.33, 1.25, DARK)
    tb(sl, title, 0.45, 0.2, 12.4, 0.9, size=29, bold=True, color=WHITE, font='Montserrat')

def notes(sl, txt):
    sl.notes_slide.notes_text_frame.text = txt

def pic(sl, path, l, t, w, h):
    return sl.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))

# ── SLIDE 1 — Title ──────────────────────────────────────────────
s1 = new_slide()
bg(s1, CREAM)
pic(s1, 'logo.png', 0.5, 0.3, 1.6, 1.6)
tb(s1, 'Mesa', 0, 1.65, 13.33, 1.25,
   size=76, bold=True, color=DARK, align=PP_ALIGN.CENTER, font='Montserrat')
tb(s1, 'Your intentions deserve a real system.', 0, 2.95, 13.33, 0.62,
   size=22, color=MID, align=PP_ALIGN.CENTER, font='Montserrat', italic=True)
rect(s1, 2.5, 3.75, 8.33, 0.05, LIGHT)
tb(s1, 'Fanny Polanco', 0, 3.92, 13.33, 0.42,
   size=17, color=CHAR, align=PP_ALIGN.CENTER)
tb(s1, 'BUS 395 — Venture Creation with AI  ·  Rollins College', 0, 4.35, 13.33, 0.38,
   size=13, color=CHAR, align=PP_ALIGN.CENTER)
tb(s1, 'Midterm Presentation  ·  June 30, 2026', 0, 4.72, 13.33, 0.38,
   size=13, color=CHAR, align=PP_ALIGN.CENTER)

notes(s1, """Before you click, introduce yourself in one sentence.

"My name is Fanny Polanco, and this is Mesa — a weekly meal-planning system for families managing health conditions like pre-diabetes on a real budget and a real schedule."

Tell the room you'll cover: the problem you confirmed in interviews, what you built, where it stands today, and what you need. Then click.

Don't read the tagline. Let it sit on screen while you speak.""")

# ── SLIDE 2 — The Problem ─────────────────────────────────────────
s2 = new_slide()
hdr(s2, 'The Problem')

rect(s2, 0.5, 1.42, 0.12, 1.5, LIGHT)
_, qtf = tb(s2, '"I\'ll do it tomorrow and never do it."', 0.75, 1.42, 11.5, 0.7,
            size=26, bold=True, color=DARK, font='Montserrat', italic=True)
ap(qtf, '— Arlene, interview subject, June 2026', size=14, color=MID, font='Lato', before=5)

_, ptf = tb(s2, 'People want to eat well for their health and their budget.', 0.5, 3.15, 12.3, 0.52,
            size=18, bold=True, color=CHAR, font='Lato')
ap(ptf, 'But there is no system connecting their grocery money, their health condition, and their weekly schedule — so the easiest default wins every time: takeout, restaurants, the same cheap processed rotation.',
   size=17, color=CHAR, font='Lato', before=8)
ap(ptf, 'The problem is not motivation. It is not knowledge. It is the complete absence of a system.',
   size=17, bold=True, color=MID, font='Lato', before=10)

tb(s2, 'Observed across 6 interviews: pre-diabetes, high blood pressure, weight management',
   0.5, 6.88, 12.3, 0.42, size=12, color=CHAR, font='Lato', italic=True)

notes(s2, """Arlene's quote names the failure mode precisely — intention exists, execution collapses. Read it once, pause, then move on. Don't explain the quote.

The core tension: everyone in your interview set knew what they should eat. Most could name the exact food they shouldn't be eating while eating it. Tere knew rice was bad for her blood sugar. Grace knew ramen was bad. They did it anyway — exhausted, no ready alternative, kids to feed.

That gap between knowing and doing is not a willpower failure. It's a systems failure. The system that should make the better choice easier than the worse choice doesn't exist yet. That's the gap Mesa closes.

If asked about sample size (6 is small): acknowledge it. The patterns were remarkably consistent, and you're continuing to collect data through the WoZ delivery and landing page.""")

# ── SLIDE 3 — What I Learned ──────────────────────────────────────
s3 = new_slide()
hdr(s3, 'What I Learned')

findings = [
    ('1  Conditions stack — the plan must capture all of them.',
     'Martha: G6PD deficiency + pre-diabetic. Cindy and her husband: weight loss + watching sugar. A plan built for one condition generates meals that conflict with the second.'),
    ('2  Intention without a system loses to the default every time.',
     'Tere, Veronica, Nelsy, Arlene — different ages, conditions, households — all defaulted to the easiest option the moment the plan broke down.'),
    ('3  The barrier is time and structure, not knowledge.',
     'Arlene: "I chose all of that because it was the easiest thing to make." And: "It takes so much time to prepare." They know what to eat. They can\'t make it happen under real conditions.'),
    ('4  No plan = the most expensive way to shop.',
     'Buying day-by-day, produce spoiling before it\'s cooked, overspending without buying healthier food. Mary spent more than she intended and came home with the same cart.'),
]

y = 1.4
for title, desc in findings:
    rect(s3, 0.5, y, 0.08, 1.05, LIGHT)
    _, ftf = tb(s3, title, 0.72, y, 12.1, 0.45, size=16, bold=True, color=DARK, font='Lato')
    ap(ftf, desc, size=14, color=CHAR, font='Lato', before=4)
    y += 1.48

notes(s3, """These four findings ARE the product architecture. Walk through each one but don't rush — finding #1 is the strongest for Q&A.

Finding 1 (conditions stack): This was new. The original problem statement assumed pre-diabetes only. Interviews revealed multi-condition households. A one-condition plan creates conflicts. This finding changed how the product collects input before generating anything.

Finding 2 (intention loses): Through-line across every interview, unprompted. Tere made rice. Veronica ordered takeout. Nelsy called it fast. Arlene said tomorrow. The product has to remove the decision point, not inform it.

Finding 3 (time/structure): Arlene is the clearest articulation. She didn't default because she forgot what to eat. She defaulted because cooking the right thing on a weeknight requires time she doesn't have. Reducing prep friction matters as much as building the plan.

Finding 4 (expensive shopping): Underused in the pitch. Your customer is already losing money — to waste, impulse buying, and unplanned day-by-day purchasing. Mesa isn't an added cost; it recovers money already being lost.""")

# ── SLIDE 4 — Solution ────────────────────────────────────────────
s4 = new_slide()
hdr(s4, 'The Solution: Mesa')

rect(s4, 0.5, 1.45, 12.33, 4.1, SAGE_BG)
_, stf = tb(s4, 'Mesa is a done-for-you weekly meal plan built around your real budget, your stores, and all of your household\'s health conditions — so the plan, not the default, decides what\'s for dinner.',
            0.85, 1.7, 11.6, 1.5, size=22, bold=True, color=DARK, font='Montserrat')
ap(stf, 'Every week: one complete grocery list, quantities matched to the whole bag — not just the recipe — and meals the whole family will actually eat. Ready before you leave for the store. Not inspiration. Not recipes. A ready-to-execute output.',
   size=18, color=CHAR, font='Lato', before=22)
ap(stf, 'The gap this fills: no product currently connects budget + perishability + schedule + blood sugar logic into one repeatable weekly system. Every competitor solves one piece. Mesa connects all four.',
   size=17, bold=True, color=MID, font='Lato', before=16)

notes(s4, """Keep this short. One clear sentence followed by two supporting points.

The sentence is on screen. Don't embellish. Let it land, then add:

"The insight that makes this different is that prior technology couldn't reason across all four constraints simultaneously — budget, perishability, schedule, and blood sugar logic. That capability only became possible with AI. That's the 'why now' for Mesa."

If asked "why not just use ChatGPT?": ChatGPT doesn't know what's on sale at your Publix this week, doesn't know your grocery budget ceiling, and doesn't remember what you made last Tuesday. The value is the integration and context, not the model itself.

Don't list features. Establish why the coordination layer matters.""")

# ── SLIDE 5 — Revenue Model & Unit Economics ──────────────────────
s5 = new_slide()
hdr(s5, 'Revenue Model & Unit Economics')

# D2C column
rect(s5, 0.5, 1.38, 5.9, 5.75, SAGE_BG)
_, d2c = tb(s5, 'DIRECT-TO-CONSUMER  ·  LIVE NOW', 0.72, 1.52, 5.5, 0.42,
            size=11, bold=True, color=MID, font='Montserrat')
ap(d2c, '$1.99 first plan · $5/month ongoing', size=23, bold=True, color=DARK, font='Montserrat', before=5)
ap(d2c, 'Live on mesa-bm3.pages.dev', size=14, color=CHAR, font='Lato', before=5)
ap(d2c, 'Testing willingness to pay now. Price shown on landing page to measure whether $1.99 creates hesitation before building payment infrastructure.',
   size=14, color=CHAR, font='Lato', before=10)
ap(d2c, 'D2C unit economics: in development.', size=15, bold=True, color=DARK, font='Lato', before=14)
ap(d2c, 'Validating willingness-to-pay before locking cost assumptions. No D2C margin stated until data supports one.',
   size=13, color=CHAR, font='Lato', before=4, italic=True)

# B2B column
rect(s5, 6.9, 1.38, 5.95, 5.75, WARM_BG)
_, b2b = tb(s5, 'B2B GROCERY CHAIN  ·  PARKED / FUTURE SCALE', 7.1, 1.52, 5.6, 0.42,
            size=11, bold=True, color=CHAR, font='Montserrat')
ap(b2b, '$5 / active user / month', size=23, bold=True, color=CHAR, font='Montserrat', before=5)
ap(b2b, 'Chain pays · end user pays nothing', size=14, color=CHAR, font='Lato', before=5)
ap(b2b, '64% gross margin  [MODELED — B2B only]', size=15, bold=True, color=CHAR, font='Lato', before=10)
ap(b2b, 'Variable cost: $0.60/user/month', size=14, color=CHAR, font='Lato', before=6)
ap(b2b, 'Fixed costs: $4,500/month', size=14, color=CHAR, font='Lato', before=4)
ap(b2b, 'Break-even: ~1,200 active users', size=14, color=CHAR, font='Lato', before=4)
ap(b2b, 'Parked: cannot reach or negotiate with a real grocery chain buyer in a semester. Testing D2C first.',
   size=13, color=CHAR, font='Lato', before=10, italic=True)

notes(s5, """Lead with D2C — that's what's live and testable today.

"Right now there's a landing page live with a real price on it — $1.99 for the first plan, $5/month after. I'm not collecting payment yet, but I'm watching whether the price makes people hesitate before filling out the form. That data will tell me whether the price is right before I invest in payment infrastructure."

On D2C unit economics: "I haven't built D2C cost assumptions yet because I don't have willingness-to-pay data yet. Stating a margin before I know if anyone pays $5/month is working backwards."

On the B2B column: "The 64% gross margin, the $4,500/month fixed costs, the 1,200 user break-even — all of that belongs to the grocery chain scenario. It's modeled against published API pricing and reasoned estimates. I'm parking it because I can't approach a real chain buyer in 12 weeks. But it's the natural scale path once D2C is validated."

If asked why the 64% margin is only in the B2B column: be precise. It's modeled, based on real Claude Sonnet API pricing plus estimates for clinical advisor, insurance, and data integration. Not guesswork, but not validated with real quotes either.""")

# ── SLIDE 6 — MVP & Landing Page ──────────────────────────────────
s6 = new_slide()
hdr(s6, 'MVP: What\'s Live')

_, url = tb(s6, 'mesa-bm3.pages.dev', 0.5, 1.42, 12.33, 0.62,
            size=28, bold=True, color=MID, font='Montserrat')
ap(url, 'Live on Cloudflare Pages · built in HTML/CSS/JS · connected to GitHub', size=13, color=CHAR, font='Lato', before=3, italic=True)

features = [
    ('Working email capture',
     'Web3Forms signup form collects name, email, and optional message. Submissions go directly to inbox. Honeypot spam protection included.'),
    ('English / Spanish toggle',
     'Full page translation on click — every line of copy, every label, every button. Built for the bilingual households in the interview set.'),
    ('Price shown to measure hesitation',
     '$1.99 intro price on the CTA, $5/month displayed as the ongoing price. Watching form completion rate as a proxy for price friction.'),
    ('Waitlist framing',
     'Positioned as "get your first plan for $1.99" — honest about current state (not yet automated) while testing real purchase intent.'),
]

y = 2.25
for title, desc in features:
    rect(s6, 0.5, y, 0.08, 0.95, LIGHT)
    _, ftf = tb(s6, title, 0.72, y, 12.1, 0.4, size=16, bold=True, color=DARK, font='Lato')
    ap(ftf, desc, size=14, color=CHAR, font='Lato', before=4)
    y += 1.13

notes(s6, """This slide is evidence of execution, not planning. The landing page is real, it's live, and it's doing something today.

The Spanish toggle tends to get a reaction. Explain why: "My interview subjects were largely Latina households. A landing page only in English signals the product wasn't built for them. The toggle took one afternoon to build and communicates something a paragraph of copy can't."

On price display: "I'm not collecting payment. But I want to know if $1.99 stops someone from completing the form. That signal is worth more right now than a payment integration I can't fulfill manually at scale."

For Q&A on why not collect payment yet: the product is still Wizard-of-Oz. Charging for something that requires 90 minutes of manual work per household is not a business — it's freelancing. First confirm willingness to pay. Then automate. Then charge.""")

# ── SLIDE 7 — Traction ────────────────────────────────────────────
s7 = new_slide()
hdr(s7, 'Traction (Honest)')

rect(s7, 0.5, 1.38, 5.9, 5.0, SAGE_BG)
_, woz = tb(s7, 'Wizard-of-Oz Delivery', 0.72, 1.52, 5.5, 0.48,
            size=17, bold=True, color=DARK, font='Montserrat')
ap(woz, '4 real households received manually built weekly meal plans:', size=15, color=CHAR, font='Lato', before=8)
ap(woz, 'Grace  ·  Martha  ·  Cindy  ·  Roberto', size=17, bold=True, color=DARK, font='Lato', before=8)
ap(woz, 'Each plan built around their actual grocery budget, health conditions, and local store sales circular. No automation yet.',
   size=14, color=CHAR, font='Lato', before=10)
ap(woz, 'Data collection in progress this week:', size=15, bold=True, color=DARK, font='Lato', before=14)
ap(woz, '· Did they shop from the plan?\n· Did they cook the meals?\n· Do they want a second week?',
   size=14, color=CHAR, font='Lato', before=6)

rect(s7, 6.9, 1.38, 5.95, 5.0, WARM_BG)
_, land = tb(s7, 'Landing Page', 7.1, 1.52, 5.6, 0.48,
             size=17, bold=True, color=DARK, font='Montserrat')
ap(land, 'mesa-bm3.pages.dev', size=16, bold=True, color=MID, font='Lato', before=8)
ap(land, 'Live and capturing sign-ups via Web3Forms.', size=15, color=CHAR, font='Lato', before=6)
ap(land, '$1.99 intro price displayed — watching form completion rate for price hesitation.',
   size=14, color=CHAR, font='Lato', before=10)
ap(land, 'English/Spanish toggle live.', size=14, color=CHAR, font='Lato', before=8)
ap(land, '⚠  Sign-up count and WoZ follow-through data will be updated before June 30 presentation.',
   size=13, bold=True, color=CHAR, font='Lato', before=16, italic=True)

notes(s7, """Be completely honest here. This is where students sometimes over-claim. Don't.

On the WoZ plans: "I manually built four meal plans for real households. Each took about 90 minutes — pulling the store circular, matching to their budget, checking against their health conditions, writing it in a format they could take to the store. That's the Wizard of Oz. The user gets the output they'd get from an automated system, but I'm the automation for now. This week I'm following up to see if any of them used it."

On data collection: be direct that you don't have results yet. "The three questions I care about: did they shop from it, did they cook from it, and do they want a second one next week. If two out of four want a second plan, I have a retention signal worth building on."

On landing page: if you have signups by June 30, say the number. If zero, say that too — and explain it's a distribution problem, not a product problem. You haven't driven traffic yet.

Professors respect honesty about what you don't know. Overclaiming in a class is a credibility risk when they ask follow-ups.""")

# ── SLIDE 8 — Pivot Log ───────────────────────────────────────────
s8 = new_slide()
hdr(s8, 'What Changed and Why')

pivots = [
    ('Week 1: ACH Reconciliation Tool',
     'Created a direct conflict of interest with my employer. Could not pursue it with integrity.',
     '→  Pivoted to Mesa'),
    ('"People have a knowledge gap"',
     'Interviews showed: they know exactly what they shouldn\'t eat. Tere knows rice. Grace knows ramen. The problem is execution under real conditions, not information.',
     '→  "People have a missing system"'),
    ('Revenue: grocery-chain-pays (B2B)',
     'Professor feedback: cannot reach or negotiate with a real chain buyer in a semester. Testing what is actually testable.',
     '→  Direct-to-consumer price test at $1.99'),
]

y = 1.42
for before, reason, after in pivots:
    rect(s8, 0.5, y, 5.5, 1.55, WARM_BG)
    _, pf = tb(s8, before, 0.7, y + 0.1, 5.1, 0.45, size=14, bold=True, color=CHAR, font='Lato')
    ap(pf, reason, size=13, color=CHAR, font='Lato', before=5, italic=True)
    tb(s8, '→', 6.1, y + 0.5, 0.6, 0.55, size=24, bold=True, color=LIGHT, align=PP_ALIGN.CENTER, font='Montserrat')
    rect(s8, 6.75, y, 6.1, 1.55, SAGE_BG)
    tb(s8, after, 6.95, y + 0.5, 5.7, 0.65, size=16, bold=True, color=DARK, font='Lato')
    y += 1.78

notes(s8, """This slide shows intellectual honesty and responsiveness — both matter in a real venture and in this class.

On the ACH pivot: "I started with a different idea entirely. Two weeks in I realized it put me in direct conflict with my employer. That's a business judgment, not a failure. I started over."

On the framing pivot: "My original problem statement said the issue was a knowledge gap. My first interview showed that was wrong. Tere knew rice was bad for her blood sugar. She ate it anyway — the only thing her kids would eat when she was exhausted at 7pm. The system gap is real. The knowledge gap I assumed doesn't exist the way I thought."

On the revenue pivot: "The original model had a grocery chain paying $5/user/month. I can't approach a chain buyer, negotiate a deal, and validate that in 12 weeks. So I tested what I could actually test: direct-to-consumer price. The landing page is live. The $1.99 is showing. I can see if people fill out the form."

Frame it clearly: pivots are not failures. They are evidence that the process is working — you're testing assumptions and updating when data disagrees.""")

# ── SLIDE 9 — What's Next ─────────────────────────────────────────
s9 = new_slide()
hdr(s9, "What's Next — Next 5 Weeks")

priorities = [
    ('1', 'Convert WoZ data into week-2 retention',
     'Grace, Martha, Cindy, and Roberto received their first plans. This week I find out whether any used it, cooked from it, and want a second one. One "yes to week 2" changes the retention assumption in the entire model.'),
    ('2', 'Drive landing page traffic · measure price hesitation at $1.99',
     'Sharing in pre-diabetes Facebook groups, WhatsApp communities, and directly with interview contacts. Watching form completion rate as a proxy for price friction. If completion drops at the $1.99 screen, I learn something before building payment infrastructure.'),
    ('3', 'Build real D2C unit economics · pressure-test B2B with one grocery contact',
     'Once I have willingness-to-pay data, I can build D2C cost assumptions. In parallel: one conversation with anyone in grocery retail — manager, loyalty tech contact, anyone — converts the B2B rate from a derived estimate to a grounded data point.'),
]

y = 1.42
for num, title, desc in priorities:
    rect(s9, 0.5, y, 0.88, 1.55, DARK)
    tb(s9, num, 0.5, y + 0.35, 0.88, 0.85, size=36, bold=True, color=WHITE,
       align=PP_ALIGN.CENTER, font='Montserrat')
    rect(s9, 1.48, y, 11.37, 1.55, SAGE_BG)
    _, ptf = tb(s9, title, 1.68, y + 0.1, 11.0, 0.48, size=16, bold=True, color=DARK, font='Lato')
    ap(ptf, desc, size=14, color=CHAR, font='Lato', before=5)
    y += 1.78

notes(s9, """Three priorities, five weeks. The professor will ask how you execute each one — be specific.

Priority 1: "The WoZ follow-up is happening this week. Three direct questions for each household: did you shop from the plan, did you cook the meals, do you want another one next week. I'll have those answers before June 30."

Priority 2: "I'm not waiting for SEO. The traffic strategy is direct outreach — posting in pre-diabetes support groups on Facebook, sharing in WhatsApp groups, reaching out to the five people I already interviewed and asking them to share the link. Small numbers, but real people who match the customer profile."

Priority 3: "On the grocery contact — I don't need a chain VP. I need one conversation with one person who works in grocery retail loyalty or health programs. LinkedIn, warm intros, anything. One real conversation replaces the most important unvalidated assumption in my B2B model."

End with confidence: "The next five weeks are about converting what I've built into evidence I can use to make real decisions about the business." """)

# ── SLIDE 10 — The Ask ────────────────────────────────────────────
s10 = new_slide()
hdr(s10, 'What I Need Help With')

asks = [
    ('A grocery retail contact',
     'Anyone with a connection inside a regional grocery chain\'s loyalty, shopper engagement, or health & wellness team. One real conversation converts the highest-risk assumption in my B2B model from a derived estimate to a data point.'),
    ('A registered dietitian or certified diabetes educator',
     'The financial model requires a clinical advisor at 20 hrs/month. I have a cost estimate ($3,000/month based on published RD rates) but no real quote. An introduction — or a reality check on whether that scope and rate are realistic — improves the model directly.'),
    ('Feedback on the $1.99 price point',
     'Is $1.99 too low to signal real value, or is it the right entry price for a budget-constrained customer? I have a hypothesis but no data yet. Experience pricing consumer health tools for this income segment would be directly useful.'),
]

y = 1.45
for title, desc in asks:
    rect(s10, 0.5, y, 0.08, 1.45, LIGHT)
    _, atf = tb(s10, title, 0.72, y, 12.1, 0.48, size=17, bold=True, color=DARK, font='Lato')
    ap(atf, desc, size=15, color=CHAR, font='Lato', before=6)
    y += 1.75

notes(s10, """End with a concrete, specific ask — not "any feedback is welcome."

"I have three specific things I need. First: a grocery retail contact. If anyone in this room knows someone inside a regional chain's loyalty or wellness program, that one introduction is worth more than any modeling I can do on my own."

"Second: a registered dietitian or CDE who might give me 30 minutes on what consulting with an early-stage health-adjacent startup looks like. I have a cost in my model but no reality check on it."

"Third: pricing feedback specifically for this customer — budget-constrained, health-managing, likely a Latina household managing on a tight grocery budget. I don't want to underprice and undermine perceived value, but I also don't want to hit a ceiling before I know where it is."

Close: "Mesa is real. The landing page is live, the plans have been delivered, and the follow-up data is coming in this week. Thank you."

Keep the close short. No recap. End with confidence.""")

# ── Save ─────────────────────────────────────────────────────────
prs.save('midterm-presentation.pptx')
print(f'Done — {len(prs.slides)} slides saved to midterm-presentation.pptx')
